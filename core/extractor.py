import fitz  # PyMuPDF
import docx
from pptx import Presentation
from pptx.util import Pt
from pathlib import Path


def extract_text(file_path: str) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(file_path)
    elif suffix in (".docx", ".doc"):
        return _extract_docx(file_path)
    elif suffix == ".txt":
        return path.read_text(encoding="utf-8", errors="ignore")
    elif suffix == ".pptx":
        return _extract_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


def _extract_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    pages = []
    for page_num, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            pages.append(f"[Page {page_num + 1}]\n{text.strip()}")
    doc.close()
    return "\n\n".join(pages)


def _extract_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        # Slide title first
        if slide.shapes.title and slide.shapes.title.text.strip():
            parts.append(f"Title: {slide.shapes.title.text.strip()}")
        # All text boxes and content shapes
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides.append(f"[Slide {i}]\n" + "\n".join(parts))
    return "\n\n".join(slides)
